# app/utils/pptx_parser.py
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from app.utils.file_ops import get_unique_image_path
import os

def parse_pptx(file_path: str):
    prs = Presentation(file_path)
    slides_data = []

    for idx, slide in enumerate(prs.slides):
        slide_content = {
            "slide_idx": idx + 1,
            "title": "",
            "paragraphs": []
        }
        if slide.shapes.title:
            slide_content["title"] = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text and text != slide_content["title"]:
                        slide_content["paragraphs"].append({
                            "text": text,
                            "level": paragraph.level
                        })
        slides_data.append(slide_content)
    return slides_data